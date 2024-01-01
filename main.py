
from flask import Flask, render_template, request, send_file
from werkzeug.utils import secure_filename
import os
from datetime import datetime, timedelta
from io import BytesIO
import base64
from docx import Document
from pptx.util import Pt, Inches
from pptx import Presentation
from pptx.dml.color import RGBColor
from PyPDF2 import PdfReader
import fitz
import requests
import chardet
import spacy
from model import *
from tempfile import NamedTemporaryFile
import uuid
import pyodbc
#import fcntl

app = Flask(__name__)
templates = "templates"

max_tokens = 900
min_tokens = 10
folder_path = "generatedfiles"

# Replace these values with your Azure SQL Database information
server = 'rajsqlserver2024.database.windows.net'
database = 'rajsqldatabase2024'
username = 'raj'
password = 'Time@123456'
driver = '{ODBC Driver 17 for SQL Server}'

def get_connection():
    return pyodbc.connect(f'DRIVER={driver};SERVER={server};DATABASE={database};UID={username};PWD={password}')

def check_connectivity():
    try:
        connection = get_connection()
        connection.close()

        print("SQL Connection successful.")
        return True
    except pyodbc.Error as e:
        print(f"Error connecting to the database: {e}")
        return False



def read_and_generate_summaries(request_id: str):
    connection = get_connection()
    cursor = connection.cursor()

    try:
        # Fetch data from tblsummary
        query = "SELECT pageno, summary FROM tblsummary WHERE requestid = ? ORDER BY pageno"
        cursor.execute(query, (request_id,))
        rows = cursor.fetchall()

        # Fetch HTML summary
        cursor.execute("EXEC usp_GetSummaries ?", (request_id,))
        html_summary = cursor.fetchone()[0]

        # Create a presentation object
        presentation = Presentation()

        # Create slides based on the data from tblsummary
        for row in rows:
            pageno, summary = row
            slide_layout = presentation.slide_layouts[5]  # Adjust layout based on your preference
            slide = presentation.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            title_shape.text = f"Page {pageno}"
            title_shape.text_frame.text = f"Page {pageno} Summary"
            title_shape.text_frame.paragraphs[0].font.size = Pt(20)
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = RGBColor(91, 155, 213)

            content_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            # Split summary into sentences
            sentences = summary.split('<br><br>')

            for i, sentence in enumerate(sentences, start=1):
                p = content_frame.add_paragraph()
                p.text = f"{sentence}"

                p.font.size = Pt(15)
                content_frame.add_paragraph().space_after = Pt(5)

        # Save the presentation to a file
        output_pptx = f"generatedfiles/{request_id}_summary.pptx"
        presentation.save(output_pptx)

        print(f"Generated File Name: {output_pptx}")
        print(f"File Exists: {os.path.exists(output_pptx)}")

        #download_link = get_pptx_download_link(output_pptx)
        # Return the file path
        success_message = "Content processed successfully"
        
        return render_template("result.html", html_summary=html_summary, request_id=request_id,success_message=success_message)

    except Exception as e:
        connection.rollback()
        error_message = f"Error retrieving or generating summaries: {str(e)}"
        print(error_message)
        return render_template("result.html", error_message=error_message)
    finally:
        connection.close()


def get_pptx_download_link(file_path):
    download_link = f'/download?file_path={file_path}'
    href = f'<a href="{download_link}" download="{os.path.basename(file_path)}">Download PowerPoint Presentation</a>'
    return href


# In your Flask app, add a route to handle the download
@app.route("/download")
def download():
    file_path = request.args.get("file_path")
    print(f"File Path: {file_path}")
    return send_file(file_path, as_attachment=True)

def delete_old_summaries():
    connection = get_connection()
    cursor = connection.cursor()

    # Call the stored procedure to delete old summaries
    cursor.execute("EXEC usp_DeleteOldSummaries")

    connection.commit()
    connection.close()


def process_text_and_display(piece_text, max_summary_length):
    summary = summarizeText(piece_text, max_summary_length)

    title = summarizeShort(piece_text)

    nlp = spacy.load("en_core_web_sm")
    sentences = [sentence.text.strip() for sentence in nlp(summary).sents]

    return sentences

def add_slide(prs, title, sentences):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(20)
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(91, 155, 213)

    content_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, sentence in enumerate(sentences, start=1):
        p = content_frame.add_paragraph()
        p.text = f"{i}. {sentence}"

        p.font.size = Pt(15)
        content_frame.add_paragraph().space_after = Pt(5)



def extract_page_pdf_text(file_bytes):
    doc = fitz.open("pdf", file_bytes)
    page_texts = []

    for page in doc:
        page_text = page.get_text()
        page_texts.append(page_text)

    doc.close()
    return page_texts

def extract_text_from_docx(file_bytes):
    doc = Document(BytesIO(file_bytes))
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"

    return text

def get_html_content(url):
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        return response.text
    except requests.exceptions.RequestException as e:
        print(f"Error retrieving HTML content from URL: {url}\nError: {str(e)}")
        return None

# Update write_summary_to_database function
def write_summary_to_database(request_id: str, page_num: int, page_content: str = None, summary: str = None):
    try:
        connection = get_connection()
        cursor = connection.cursor()

        cursor.execute('''
            EXEC usp_InsertSummaryData @RequestID=?, @PageNo=?, @PageContent=?, @Summary=?
        ''', (request_id, page_num, page_content, summary))

        connection.commit()
    except pyodbc.Error as e:
        # Handle specific database-related errors
        print(f"Database error: {e}")
        # Optionally, log the error using a logging framework

    except Exception as e:
        # Handle other types of exceptions
        print(f"An error occurred: {e}")
        # Optionally, log the error using a logging framework

    finally:
        # Close the connection in a finally block to ensure it happens regardless of exceptions
        if connection:
            connection.close()



@app.route('/', methods=['GET', 'POST'])
def index():
    return render_template("index.html")

def generate_request_id():
    return str(uuid.uuid4())

@app.route('/generate_summary', methods=['POST'])
def generate_summary():
    try:
        if check_connectivity():
            request_id = generate_request_id()
            delete_old_summaries()  # Delete old summaries before processing

            if request.form.get("input_choice") == "Upload File":
                uploaded_file = request.files['uploaded_file']
                file_bytes = uploaded_file.read()
                file_extension = os.path.splitext(uploaded_file.filename)[-1].lower()

                if file_extension == ".pdf":
                    page_texts = extract_page_pdf_text(file_bytes)
                    
                    print(f'Raj: Total Pages: {len(page_texts)}')

                    for page_num, page_text in enumerate(page_texts, start=1):
                        nested_sentences = create_nested_sentences(page_text, token_max_length=900)
                        summary = ""
                        print(f'PageNum: {page_num}')
                        for idx, nested in enumerate(nested_sentences):
                            concatenated_text = " ".join(nested)
                            sentences = process_text_and_display(concatenated_text, max_tokens)

                            for i, sentence in enumerate(sentences, start=1):
                                summary += f"  {i}. {sentence} <br>"

                        write_summary_to_database(request_id, page_num, page_text, summary)

                elif file_extension == ".docx":
                    page_text = extract_text_from_docx(file_bytes)
                    nested_sentences = create_nested_sentences(page_text, token_max_length=max_tokens)
                    summary = ""

                    for idx, nested in enumerate(nested_sentences):
                        concatenated_text = " ".join(nested)
                        sentences = process_text_and_display(concatenated_text, max_tokens)

                        for i, sentence in enumerate(sentences, start=1):
                            summary += f"{i}. {sentence}<br>"

                    write_summary_to_database(request_id, 1, page_text, summary)  # Assuming single page for DOCX

                elif file_extension == ".txt":
                    file_encoding = chardet.detect(file_bytes)['encoding'] or 'utf-8'

                    text = file_bytes.decode(file_encoding)
                    nested_sentences = create_nested_sentences(text, token_max_length=max_tokens)
                    summary = ""

                    for idx, nested in enumerate(nested_sentences):
                        concatenated_text = " ".join(nested)
                        sentences = process_text_and_display(concatenated_text, max_tokens)

                        for i, sentence in enumerate(sentences, start=1):
                            summary += f"{i}. {sentence}<br>"

                        write_summary_to_database(request_id, 1, concatenated_text, summary)  # Assuming single page for TXT

                else:
                    raise ValueError("Unsupported file format")

            elif request.form.get("input_choice") == "Paste Text":
                pasted_text = request.form.get("pasted_text")
                nested_sentences = create_nested_sentences(pasted_text, token_max_length=max_tokens)
                summary = ""

                for idx, nested in enumerate(nested_sentences):
                    concatenated_text = " ".join(nested)
                    sentences = process_text_and_display(concatenated_text, max_tokens)

                    for i, sentence in enumerate(sentences, start=1):
                        summary += f"  {i}. {sentence}<br><br>"

                write_summary_to_database(request_id, 1, pasted_text, summary)  # Assuming single page for pasted text

            # Retrieve summaries from the database
            
            #read_summaries_from_databasePPTX(request_id)
            return read_and_generate_summaries(request_id)

            #summaries = read_summaries_from_database(request_id)

            # if summaries:
            #     return render_template(
            #         "result.html", summaries=summaries, success_message="Content processed successfully"
            #     )
            # else:
            #     error_message = "No summaries found for the given request ID."
            #     return render_template("result.html", error_message=error_message)
                
            # Retrieve summaries from the database and render template
            # return render_template(
            #     "result.html",
            #     summaries=read_summaries_from_database(request_id),
            #     success_message="Content processed successfully"
            # ) if read_summaries_from_database(request_id) else render_template(
            #     "result.html",
            #     error_message="No summaries found for the given request ID."
            # )


    except Exception as e:
        error_message = f"Error: {str(e)}"
        return render_template("result.html", error_message=error_message)

if __name__ == '__main__':
    app.run(debug=True)