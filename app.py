from flask import Flask, render_template, request, send_file
from datetime import datetime, timedelta
import base64
from docx import Document
from io import BytesIO
import os
from io import StringIO
from datetime import datetime
import io
from flask import send_file
#from flask_ngrok import run_with_ngrok
import spacy
from bs4 import BeautifulSoup
from PyPDF2 import PdfReader
import validators
import fitz
import requests
#import torch
import re
import chardet
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx import Presentation
from pptx.util import Inches
from flask import send_file

# Import your functions from model.py
from model import *

app = Flask(__name__)
#run_with_ngrok(app)  # Start ngrok when app is run

# Function to delete files older than a specified time
def delete_old_files(folder_path, time_threshold_seconds=120):
    current_time = datetime.now()

    for filename in os.listdir(folder_path):
        if filename.startswith("Text_Summary_Presentation_") and filename.endswith(".pptx"):
            file_path = os.path.join(folder_path, filename)

            # Get the file creation timestamp
            created_timestamp = datetime.fromtimestamp(os.path.getctime(file_path))

            # Calculate the time difference
            time_difference = current_time - created_timestamp

            # If the file is older than the threshold, delete it
            if time_difference.total_seconds() > time_threshold_seconds:
                os.remove(file_path)

# Function to generate a downloadable link for the PowerPoint presentation
def get_pptx_download_link(file_path):
    with open(file_path, 'rb') as f:
        pptx_data = f.read()
    encoded_pptx = base64.b64encode(pptx_data).decode()
    href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{encoded_pptx}" download="Text_Summary_Presentation.pptx">Download PowerPoint Presentation</a>'
    return href

def add_slide(prs, title, sentences):
    slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title with background color
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(20)  # Set font size for the title
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(91, 155, 213)  # Set background color for the title

    # Add content
    content_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, sentence in enumerate(sentences, start=1):
        p = content_frame.add_paragraph()
        p.text = f"{i}. {sentence}"

        p.font.size = Pt(15)  # Set font size for the content
        # Add an empty line after each sentence
        content_frame.add_paragraph().space_after = Pt(5)

def clean_text(web_text):
    return ' '.join(line.strip() for line in web_text.splitlines() if line.strip())

# Create PowerPoint presentation
#presentation = Presentation()

def process_text_and_display(piece_text, max_summary_length, presentation):
    summary = summarizeText(piece_text, max_summary_length)
    title = summarizeShort(piece_text)

    # Use spaCy for sentence tokenization
    nlp = spacy.load("en_core_web_sm")
    sentences = [sentence.text.strip() for sentence in nlp(summary).sents]

    add_slide(presentation, title, sentences)

def get_html_content(url):
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()  # Raises an HTTPError for bad responses
        return response.text
    except requests.exceptions.RequestException as e:
        return None

def process_url(url, max_summary_length):
    if not validators.url(url):
        return None

    html_content = get_html_content(url)

    if html_content:
        nlp = spacy.load("en_core_web_sm", disable=["parser", "ner"])
        soup = BeautifulSoup(html_content, "html.parser")
        title = soup.title.string
        web_text = soup.get_text()

        return {
            'url': url,
            'title': title,
            'cleaned_text': web_text
        }

    else:
        return None

def extract_page_pdf_text(file_bytes):
    doc = fitz.open("pdf", file_bytes)
    page_texts = []

    for page in doc:
        page_text = page.get_text()
        page_texts.append(page_text)

    doc.close()
    return page_texts

def extract_text_from_docx(file_bytes):
    doc = Document(BytesIO(file_bytes))
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"

    return text

def clean_text_for_summarization(text):
    cleaned_text = ' '.join(text.split())
    return cleaned_text

def extract_text_from_file(file_bytes, file_extension):
    if file_extension == ".pdf":
        return extract_page_pdf_text(file_bytes)
    elif file_extension == ".docx":
        return extract_text_from_docx(file_bytes)
    else:
        return None

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate_summary', methods=['POST'])
def generate_summary():
    max_summary_length = int(request.form['max_summary_length'])
    input_choice = request.form['input_choice']

    if input_choice == "Website URL":

        presentation = Presentation()
        url_input = request.form['url_input']

        if url_input:
            # Call the function to delete old files
            delete_old_files("docs")

            result = process_url(url_input, max_summary_length)
            if result is not None:
                title = result['title']
                web_text = result['cleaned_text']
                nested_sentences = create_nested_sentences(web_text, token_max_length=900)

                for idx, nested in enumerate(nested_sentences):
                    concatenated_text = " ".join(nested)
                    process_text_and_display(concatenated_text, max_summary_length,presentation)

                timestamp = datetime.now().strftime("%Y-%m-%d-%H-%M-%S")
                pptx_filename = f"Text_Summary_Presentation_{timestamp}.pptx"
                pptx_filepath = os.path.join("docs", pptx_filename)
                presentation.save(pptx_filepath)

                return send_file(pptx_filepath, as_attachment=True)
    elif input_choice == "Upload File":

        presentation = Presentation()
        uploaded_file = request.files['uploaded_file']

        if uploaded_file:
            # Call the function to delete old files
            delete_old_files("docs")

            file_bytes = uploaded_file.read()
            file_extension = os.path.splitext(uploaded_file.filename)[-1].lower()

            if file_extension == ".pdf":
                page_texts = extract_page_pdf_text(file_bytes)

                for page_num, page_text in enumerate(page_texts, start=1):
                    nested_sentences = create_nested_sentences(page_text, token_max_length=900)
                    for idx, nested in enumerate(nested_sentences):
                        concatenated_text = " ".join(nested)
                        process_text_and_display(concatenated_text, max_summary_length,presentation)

            elif file_extension == ".docx":
                page_text = extract_text_from_docx(file_bytes)

                nested_sentences = create_nested_sentences(page_text, token_max_length=900)

                for idx, nested in enumerate(nested_sentences):
                    concatenated_text = " ".join(nested)
                    process_text_and_display(concatenated_text, max_summary_length,presentation)

            elif file_extension == ".txt":
                text = file_bytes.decode('utf-8')
                nested_sentences = create_nested_sentences(text, token_max_length=900)

                for idx, nested in enumerate(nested_sentences):
                    concatenated_text = " ".join(nested)
                    process_text_and_display(concatenated_text, max_summary_length,presentation)

            timestamp = datetime.now().strftime("%Y-%m-%d-%H-%M-%S")
            pptx_filename = f"Text_Summary_Presentation_{timestamp}.pptx"
            pptx_filepath = os.path.join("docs", pptx_filename)
            presentation.save(pptx_filepath)

            return send_file(pptx_filepath, as_attachment=True)

    elif input_choice == "Paste Text":

        presentation = Presentation()
        pasted_text = request.form['pasted_text']

        if pasted_text:
            # Call the function to delete old files
            delete_old_files("docs")

            nested_sentences = create_nested_sentences(pasted_text, token_max_length=900)

            for idx, nested in enumerate(nested_sentences):
                concatenated_text = " ".join(nested)
                process_text_and_display(concatenated_text, max_summary_length,presentation)

            timestamp = datetime.now().strftime("%Y-%m-%d-%H-%M-%S")
            pptx_filename = f"Text_Summary_Presentation_{timestamp}.pptx"
            pptx_filepath = os.path.join("docs", pptx_filename)
            presentation.save(pptx_filepath)

            return send_file(pptx_filepath, as_attachment=True)

    return "Invalid input choice or missing data."

if __name__ == '__main__':
    app.run(debug=True)
