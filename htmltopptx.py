from flask import Flask, render_template, request
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import pyodbc
import base64
import os
import uuid

app = Flask(__name__)

# Replace these values with your Azure SQL Database information
server = 'rajsqlserver2024.database.windows.net'
database = 'rajsqldatabase2024'
username = 'raj'
password = 'Time@123456'
driver = '{ODBC Driver 17 for SQL Server}'


    
def get_connection():
    return pyodbc.connect(f'DRIVER={driver};SERVER={server};DATABASE={database};UID={username};PWD={password}')

def get_pptx_download_link(file_path):
    with open(file_path, 'rb') as f:
        pptx_data = f.read()
    encoded_pptx = base64.b64encode(pptx_data).decode()
    href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{encoded_pptx}" download="Text_Summary_Presentation.pptx">Download PowerPoint Presentation</a>'
    return href

def read_summaries_from_database(request_id: str):
    connection = get_connection()
    cursor = connection.cursor()

    try:
        # Fetch data from tblsummary
        cursor.execute("SELECT pageno, summary FROM tblsummary WHERE request_id = ?", (request_id,))
        rows = cursor.fetchall()

        # Create a presentation object
        presentation = Presentation()

        # Create slides based on the data from tblsummary
        for row in rows:
            pageno, summary = row
            slide_layout = presentation.slide_layouts[5]  # Adjust layout based on your preference
            slide = presentation.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            title_shape.text = f"Page {pageno} Summary"
            title_shape.text_frame.text = f"Page {pageno} Summary"
            title_shape.text_frame.paragraphs[0].font.size = Pt(20)
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = RGBColor(91, 155, 213)

            content_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            # Split summary into sentences
            sentences = summary.split('<br>')

            for i, sentence in enumerate(sentences, start=1):
                p = content_frame.add_paragraph()
                p.text = f"{i}. {sentence}"

                p.font.size = Pt(15)
                content_frame.add_paragraph().space_after = Pt(5)

        # Save the presentation to a file
        output_pptx = f"generatedfiles/{request_id}_summary.pptx"
        presentation.save(output_pptx)

        # Return the file path or a download link
        download_link = get_pptx_download_link(output_pptx)
        return render_template("result.html", download_link=download_link)

    except Exception as e:
        connection.rollback()
        error_message = f"Error retrieving summaries: {str(e)}"
        return render_template("result.html", error_message=error_message)
    finally:
        connection.close()

# Add your other functions and routes here...



if __name__ == '__main__':
    app.run(debug=True)
