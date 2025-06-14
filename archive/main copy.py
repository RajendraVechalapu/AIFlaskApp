
from fastapi import Depends, FastAPI, Form, File, UploadFile, Request
from fastapi.responses import HTMLResponse
from fastapi.templating import Jinja2Templates
import os
from datetime import datetime, timedelta
from io import BytesIO
import base64
from docx import Document
from pptx.util import Pt, Inches
from pptx import Presentation
from pptx.dml.color import RGBColor
from PyPDF2 import PdfReader
import fitz
import requests
import chardet
import spacy
from model import *
from tempfile import NamedTemporaryFile
import uuid
import pyodbc

app = FastAPI()
templates = Jinja2Templates(directory="templates")

max_tokens = 900
min_tokens = 10
folder_path = "generatedfiles"

# Replace these values with your Azure SQL Database information
server = 'rajsqlserver2024.database.windows.net'
database = 'rajsqldatabase2024'
username = 'raj'
password = 'REMOVED_PASSWORD'
driver = '{ODBC Driver 17 for SQL Server}'

def get_connection():
    return pyodbc.connect(f'DRIVER={driver};SERVER={server};DATABASE={database};UID={username};PWD={password}')

def check_connectivity():
    try:
        connection = get_connection()
        connection.close()

        print("SQL Connection successful.")
        return True
    except pyodbc.Error as e:
        print(f"Error connecting to the database: {e}")
        return False



def read_summaries_from_database(request_id: str):
    connection = get_connection()
    cursor = connection.cursor()
    cursor.execute('SELECT PageNo, Summary, CreatedDateTime FROM tblSummary WHERE RequestID = ?', (request_id,))
    rows = cursor.fetchall()
    connection.close()

    summaries = []
    for row in rows:
        PageNo, summary, createddatetime = row
        summaries.append({"pageno": PageNo, "summary": summary, "createddatetime": createddatetime})

    return summaries

def delete_old_summaries():
    connection = get_connection()
    cursor = connection.cursor()

    # Call the stored procedure to delete old summaries
    cursor.execute("EXEC usp_DeleteOldSummaries")

    connection.commit()
    connection.close()


def process_text_and_display(piece_text, max_summary_length):
    summary = summarizeText(piece_text, max_summary_length)

    title = summarizeShort(piece_text)

    nlp = spacy.load("en_core_web_sm")
    sentences = [sentence.text.strip() for sentence in nlp(summary).sents]

    return sentences

def add_slide(prs, title, sentences):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(20)
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(91, 155, 213)

    content_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(2), width=Inches(8), height=Inches(4))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for i, sentence in enumerate(sentences, start=1):
        p = content_frame.add_paragraph()
        p.text = f"{i}. {sentence}"

        p.font.size = Pt(15)
        content_frame.add_paragraph().space_after = Pt(5)

def get_pptx_download_link(file_path):
    with open(file_path, 'rb') as f:
        pptx_data = f.read()
    encoded_pptx = base64.b64encode(pptx_data).decode()
    href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{encoded_pptx}" download="Text_Summary_Presentation.pptx">Download PowerPoint Presentation</a>'
    return href

def extract_page_pdf_text(file_bytes):
    doc = fitz.open("pdf", file_bytes)
    page_texts = []

    for page in doc:
        page_text = page.get_text()
        page_texts.append(page_text)

    doc.close()
    return page_texts

def extract_text_from_docx(file_bytes):
    doc = Document(BytesIO(file_bytes))
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"

    return text

def get_html_content(url):
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        return response.text
    except requests.exceptions.RequestException as e:
        print(f"Error retrieving HTML content from URL: {url}\nError: {str(e)}")
        return None

# Update write_summary_to_database function
def write_summary_to_database(request_id: str, page_num: int, page_content: str = None, summary: str = None):
    try:
        connection = get_connection()
        cursor = connection.cursor()

        cursor.execute('''
            EXEC usp_InsertSummaryData @RequestID=?, @PageNo=?, @PageContent=?, @Summary=?
        ''', (request_id, page_num, page_content, summary))

        connection.commit()
    except pyodbc.Error as e:
        # Handle specific database-related errors
        print(f"Database error: {e}")
        # Optionally, log the error using a logging framework

    except Exception as e:
        # Handle other types of exceptions
        print(f"An error occurred: {e}")
        # Optionally, log the error using a logging framework

    finally:
        # Close the connection in a finally block to ensure it happens regardless of exceptions
        if connection:
            connection.close()



@app.get('/')
def index(request: Request):
    #print("Vechalapu Rajendra Simhadri Appala Naidu.")
    return templates.TemplateResponse("index.html", {"request": request})

def generate_request_id():
    return str(uuid.uuid4())
@app.post('/generate_summary')
async def generate_summary(
    request: Request,
    max_summary_length: int = Form(...),
    input_choice: str = Form(None),
    uploaded_file: UploadFile = File(...),
    pasted_text: str = Form(None),
    request_id: str = Depends(generate_request_id),
):
    try:
        if check_connectivity():
            delete_old_summaries()  # Delete old summaries before processing

            if input_choice == "Upload File":
                file_bytes = await uploaded_file.read()
                file_extension = os.path.splitext(uploaded_file.filename)[-1].lower()

                if file_extension == ".pdf":
                    page_texts = extract_page_pdf_text(file_bytes)
                    
                    print(f'Raj: Total Pages: {len(page_texts)}')

                    for page_num, page_text in enumerate(page_texts, start=1):
                        nested_sentences = create_nested_sentences(page_text, token_max_length=900)
                        summary = ""
                        print(f'PageNum: {page_num}')
                        for idx, nested in enumerate(nested_sentences):
                            concatenated_text = " ".join(nested)
                            sentences = process_text_and_display(concatenated_text, max_summary_length)

                            for i, sentence in enumerate(sentences, start=1):
                                summary += f"  {i}. {sentence} <br>"

                        write_summary_to_database(request_id,page_num, page_text, summary)

                elif file_extension == ".docx":
                    page_text = extract_text_from_docx(file_bytes)
                    nested_sentences = create_nested_sentences(page_text, token_max_length=max_tokens)
                    summary = ""

                    for idx, nested in enumerate(nested_sentences):
                        concatenated_text = " ".join(nested)
                        sentences = process_text_and_display(concatenated_text, max_summary_length)

                        for i, sentence in enumerate(sentences, start=1):
                            summary += f"{i}. {sentence}<br>"

                    write_summary_to_database(request_id, 1, page_text, summary)  # Assuming single page for DOCX

                elif file_extension == ".txt":
                    file_encoding = chardet.detect(file_bytes)['encoding'] or 'utf-8'

                    text = file_bytes.decode(file_encoding)
                    nested_sentences = create_nested_sentences(text, token_max_length=max_tokens)
                    summary = ""

                    for idx, nested in enumerate(nested_sentences):
                        concatenated_text = " ".join(nested)
                        sentences = process_text_and_display(concatenated_text, max_summary_length)

                        for i, sentence in enumerate(sentences, start=1):
                            summary += f"{i}. {sentence}<br>"

                    write_summary_to_database(request_id, 1, text, summary)  # Assuming single page for TXT

                else:
                    raise ValueError("Unsupported file format")

                # Write summary to database
                #write_summary_to_database(request_id, "", summary)

            elif input_choice == "Paste Text":
                nested_sentences = create_nested_sentences(pasted_text, token_max_length=max_tokens)
                summary = ""

                for idx, nested in enumerate(nested_sentences):
                    concatenated_text = " ".join(nested)
                    sentences = process_text_and_display(concatenated_text, max_summary_length)

                    for i, sentence in enumerate(sentences, start=1):
                        summary += f"  {i}. {sentence}<br><br>"

                write_summary_to_database(request_id, 1, pasted_text, summary)  # Assuming single page for pasted text


                # Write summary to database
                #write_summary_to_database(request_id, "", summary)

            # Retrieve summary from database
            # Retrieve summaries from the database
            summaries = read_summaries_from_database(request_id)

            if summaries:
                return templates.TemplateResponse(
                    "result.html", {"request": request, "summaries": summaries, "success_message": "Content processed successfully"}
                )
            else:
                error_message = "No summaries found for the given request ID."
                return templates.TemplateResponse(
                    "result.html", {"request": request, "error_message": error_message}
                )

    except Exception as e:
        error_message = f"Error: {str(e)}"
        return templates.TemplateResponse(
            "result.html", {"request": request, "error_message": error_message}
        )
